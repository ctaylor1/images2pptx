import logging
import sys
import os
import yaml

from pptx import Presentation
from pptx.util import Inches, Pt
from pytesseract import image_to_string
from PIL import Image


def setup_logger():
    """
    Configures the root logger to output messages to both stdout and a file in ./logs/app.log.
    """
    log_directory = "logs"
    os.makedirs(log_directory, exist_ok=True)
    log_file = os.path.join(log_directory, "app.log")

    logging.basicConfig(
        level=logging.INFO,  # Change to logging.DEBUG for more verbosity
        format='%(asctime)s [%(levelname)s] %(message)s',
        handlers=[
            logging.StreamHandler(sys.stdout),
            logging.FileHandler(log_file, mode='a', encoding='utf-8')
        ]
    )


def load_config(config_file: str) -> dict:
    """
    Loads configuration from a YAML file.
    
    :param config_file: Path to the YAML config file.
    :return: A dictionary with configuration data.
    """
    if not os.path.isfile(config_file):
        logging.error(f"Config file '{config_file}' not found.")
        sys.exit(1)

    try:
        with open(config_file, 'r', encoding='utf-8') as file:
            config = yaml.safe_load(file)
            if config is None:
                raise ValueError("Empty or invalid YAML structure.")
    except (yaml.YAMLError, ValueError) as e:
        logging.error(f"Failed to parse YAML config: {e}")
        sys.exit(1)

    return config


def validate_config(config: dict) -> None:
    """
    Validates the essential keys and types in the configuration,
    logs errors and exits if required keys are missing or invalid.
    """
    # Check top-level keys
    if "paths" not in config:
        logging.error("Missing 'paths' section in config.yaml.")
        sys.exit(1)
    if "presentation" not in config:
        logging.error("Missing 'presentation' section in config.yaml.")
        sys.exit(1)

    # Validate 'paths'
    paths = config["paths"]
    required_paths_keys = ["images_folder", "output_folder", "output_filename"]
    for key in required_paths_keys:
        if key not in paths:
            logging.error(f"Missing '{key}' under 'paths' in config.yaml.")
            sys.exit(1)

    images_folder = paths["images_folder"]
    output_folder = paths["output_folder"]
    output_filename = paths["output_filename"]

    # Check that images_folder exists
    if not os.path.isdir(images_folder):
        logging.error(f"Images folder '{images_folder}' does not exist or is not a directory.")
        sys.exit(1)

    if not isinstance(output_filename, str) or not output_filename.endswith(".pptx"):
        logging.error(f"Output filename '{output_filename}' must be a string ending with .pptx.")
        sys.exit(1)

    # Validate 'presentation'
    presentation = config["presentation"]
    # NEW / CHANGED: No longer require slide_width_inches or slide_height_inches
    # But we do expect slide_size_option (fallback to 'widescreen' if invalid).
    required_presentation_keys = [
        "textbox_left_inches",
        "textbox_top_inches",
        "textbox_width_inches",
        "textbox_height_inches",
        "image_left_inches",
        "image_top_inches",
        "image_scale_percent",
        "text_font_size"
    ]
    # We'll treat slide_size_option as optional (defaults to "widescreen")
    for key in required_presentation_keys:
        if key not in presentation:
            logging.error(f"Missing '{key}' under 'presentation' in config.yaml.")
            sys.exit(1)
        try:
            float(presentation[key])  # Just ensure it can be cast to float
        except (TypeError, ValueError):
            logging.error(f"'{key}' under 'presentation' must be a numerical value.")
            sys.exit(1)

    # Validate 'extensions' if present
    if "extensions" in config:
        if not isinstance(config["extensions"], list):
            logging.error("'extensions' in config.yaml should be a list of file extensions.")
            sys.exit(1)
        for ext in config["extensions"]:
            if not (isinstance(ext, str) and ext.startswith(".")):
                logging.error(f"Invalid extension '{ext}' in 'extensions'; must be a string like '.png'.")
                sys.exit(1)


def create_powerpoint_slides(config: dict):
    """
    Reads image files from a folder, performs OCR, and creates a PowerPoint
    presentation with the images and extracted text.

    :param config: Configuration dictionary loaded from config.yaml
    """
    # Extract config data
    images_folder = config["paths"]["images_folder"]
    output_folder = config["paths"]["output_folder"]
    output_filename = config["paths"]["output_filename"]

    presentation_cfg = config["presentation"]

    textbox_left = float(presentation_cfg["textbox_left_inches"])
    textbox_top = float(presentation_cfg["textbox_top_inches"])
    textbox_width = float(presentation_cfg["textbox_width_inches"])
    textbox_height = float(presentation_cfg["textbox_height_inches"])

    image_left = float(presentation_cfg["image_left_inches"])
    image_top = float(presentation_cfg["image_top_inches"])
    image_scale_percent = float(presentation_cfg["image_scale_percent"])
    text_font_size = float(presentation_cfg["text_font_size"])

    # NEW / CHANGED: Slide Size Option
    slide_size_option = presentation_cfg.get("slide_size_option", "widescreen").lower()
    # Map of recognized slide size options (width_in, height_in)
    size_map = {
        "standard":   (10.0, 7.5),     # 4:3 ratio
        "widescreen": (13.3333, 7.5)   # 16:9 ratio
    }
    # Fallback to widescreen if invalid
    if slide_size_option not in size_map:
        logging.warning(f"Invalid slide_size_option '{slide_size_option}' provided. Defaulting to 'widescreen'.")
        slide_size_option = "widescreen"
    slide_width, slide_height = size_map[slide_size_option]

    # Get allowed extensions from config (fallback to .png if none specified)
    allowed_extensions = config.get("extensions", [".png"])

    # Log folder paths and file output
    logging.info(f"Images folder: {images_folder}")
    logging.info(f"Output folder: {output_folder}")
    logging.info(f"Output filename: {output_filename}")
    logging.info(f"Allowed extensions: {allowed_extensions}")
    logging.info(f"Using slide size option: {slide_size_option} ({slide_width}in x {slide_height}in)")

    # Ensure the output directory exists or create it
    if not os.path.exists(output_folder):
        logging.info(f"Output folder '{output_folder}' does not exist. Creating it.")
        try:
            os.makedirs(output_folder, exist_ok=True)
        except OSError as e:
            logging.error(f"Failed to create output folder '{output_folder}': {e}")
            sys.exit(1)

    # Build the full output path
    full_output_path = os.path.join(output_folder, output_filename)

    # Create a PowerPoint presentation
    presentation = Presentation()
    # Set the slide dimensions
    presentation.slide_width = Inches(slide_width)
    presentation.slide_height = Inches(slide_height)

    # Sort files in alphabetical order, filter by allowed extensions
    try:
        all_files = sorted(os.listdir(images_folder))
    except OSError as e:
        logging.error(f"Error reading images folder '{images_folder}': {e}")
        sys.exit(1)

    image_files = [
        f for f in all_files
        if any(f.lower().endswith(ext) for ext in allowed_extensions)
    ]

    if not image_files:
        logging.warning("No valid image files found in the specified folder.")

    for image_file in image_files:
        image_path = os.path.join(images_folder, image_file)
        logging.info(f"Processing image: {image_path}")

        # Open image and run OCR
        try:
            with Image.open(image_path) as img:
                img.seek(0)  # For multi-frame images
                extracted_text = image_to_string(img)

                # Keep aspect ratio when scaling
                orig_width_px, orig_height_px = img.size
        except OSError as e:
            logging.error(f"Could not open or read image '{image_path}': {e}")
            continue

        # Add a blank slide
        slide_layout_index = 6  # Typically a "Blank" slide layout
        slide = presentation.slides.add_slide(presentation.slide_layouts[slide_layout_index])

        # Calculate scaled dimensions in inches
        dpi_assumption = 96.0
        base_width_in = orig_width_px / dpi_assumption
        base_height_in = orig_height_px / dpi_assumption
        scale_factor = image_scale_percent / 100.0
        scaled_width_in = base_width_in * scale_factor
        scaled_height_in = base_height_in * scale_factor

        # Insert the image
        try:
            slide.shapes.add_picture(
                image_path,
                Inches(image_left),
                Inches(image_top),
                width=Inches(scaled_width_in),
                height=Inches(scaled_height_in)
            )
        except OSError as e:
            logging.error(f"Could not add image '{image_path}' to slide: {e}")
            continue

        # Add the extracted text box
        text_box = slide.shapes.add_textbox(
            Inches(textbox_left),
            Inches(textbox_top),
            Inches(textbox_width),
            Inches(textbox_height)
        )
        text_frame = text_box.text_frame
        text_frame.text = extracted_text

        # Set font size for all text within the text box
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(text_font_size)

    # Save the PowerPoint
    try:
        presentation.save(full_output_path)
        logging.info(f"PowerPoint presentation saved to: {full_output_path}")
    except OSError as e:
        logging.error(f"Failed to save PowerPoint to '{full_output_path}': {e}")
        sys.exit(1)


def main():
    setup_logger()
    config = load_config("config.yaml")
    validate_config(config)
    create_powerpoint_slides(config)


if __name__ == "__main__":
    main()